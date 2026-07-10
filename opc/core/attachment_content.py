"""Attachment content helpers for extraction and multimodal routing."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Iterable

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".yaml", ".yml",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css",
    ".xml", ".toml", ".ini", ".cfg", ".log", ".sql",
}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}


class _PreviewAccumulator:
    def __init__(self, max_chars: int) -> None:
        self.max_chars = max(0, max_chars)
        self.parts: list[str] = []
        self.used = 0

    def add(self, text: str) -> bool:
        normalized = str(text or "").strip()
        if not normalized:
            return False
        if self.used >= self.max_chars:
            return True

        available = self.max_chars - self.used
        if len(normalized) > available:
            normalized = normalized[:available].rstrip()
        if not normalized:
            return True

        self.parts.append(normalized)
        self.used += len(normalized) + 1
        return self.used >= self.max_chars

    def render(self) -> str:
        return "\n".join(self.parts).strip()


def attachment_suffix(filename: str) -> str:
    return Path(filename).suffix.lower()


def is_text_like_attachment(filename: str, mime_type: str) -> bool:
    if mime_type.startswith("text/"):
        return True
    return attachment_suffix(filename) in TEXT_EXTENSIONS


def can_extract_text(filename: str, mime_type: str) -> bool:
    return is_text_like_attachment(filename, mime_type) or attachment_suffix(filename) in OFFICE_EXTENSIONS


def extract_attachment_text(
    filename: str,
    mime_type: str,
    raw: bytes,
    *,
    max_chars: int = 4000,
) -> str:
    suffix = attachment_suffix(filename)
    if is_text_like_attachment(filename, mime_type):
        return _clip_text(raw.decode("utf-8", errors="replace").strip(), max_chars)
    if suffix == ".docx":
        return _extract_docx_text(raw, max_chars=max_chars)
    if suffix == ".xlsx":
        return _extract_xlsx_text(raw, max_chars=max_chars)
    if suffix == ".pptx":
        return _extract_pptx_text(raw, max_chars=max_chars)
    return ""


def _clip_text(text: str, max_chars: int) -> str:
    text = str(text or "").strip()
    if len(text) <= max_chars:
        return text
    return f"{text[:max_chars].rstrip()}\n...[truncated]"


def _extract_docx_text(raw: bytes, *, max_chars: int) -> str:
    from docx import Document

    acc = _PreviewAccumulator(max_chars)
    doc = Document(BytesIO(raw))

    for para in doc.paragraphs:
        if acc.add(para.text):
            return _clip_text(acc.render(), max_chars)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = " | ".join(cell for cell in cells if cell)
            if acc.add(line):
                return _clip_text(acc.render(), max_chars)

    return _clip_text(acc.render(), max_chars)


def _extract_xlsx_text(raw: bytes, *, max_chars: int) -> str:
    from openpyxl import load_workbook

    acc = _PreviewAccumulator(max_chars)
    workbook = load_workbook(BytesIO(raw), read_only=True, data_only=True)
    try:
        for sheet in workbook.worksheets[:5]:
            if acc.add(f"# Sheet: {sheet.title}"):
                break
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                values = [_normalize_excel_cell(value) for value in row[:16]]
                if not any(values):
                    continue
                row_count += 1
                if acc.add("\t".join(values)):
                    return _clip_text(acc.render(), max_chars)
                if row_count >= 80:
                    break
    finally:
        workbook.close()

    return _clip_text(acc.render(), max_chars)


def _normalize_excel_cell(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, float):
        text = f"{value:.6f}".rstrip("0").rstrip(".")
        return text or "0"
    return str(value).strip()


def _extract_pptx_text(raw: bytes, *, max_chars: int) -> str:
    from pptx import Presentation

    acc = _PreviewAccumulator(max_chars)
    presentation = Presentation(BytesIO(raw))

    for index, slide in enumerate(list(presentation.slides)[:20], start=1):
        if acc.add(f"# Slide {index}"):
            break
        for text in _iter_slide_text(slide.shapes):
            if acc.add(text):
                return _clip_text(acc.render(), max_chars)

    return _clip_text(acc.render(), max_chars)


def _iter_slide_text(shapes: Iterable[object]) -> Iterable[str]:
    for shape in shapes:
        text = getattr(shape, "text", "")
        if isinstance(text, str) and text.strip():
            yield text

        table = getattr(shape, "table", None)
        if table is not None:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                line = " | ".join(cell for cell in cells if cell)
                if line:
                    yield line

        subshapes = getattr(shape, "shapes", None)
        if subshapes is not None:
            yield from _iter_slide_text(subshapes)
